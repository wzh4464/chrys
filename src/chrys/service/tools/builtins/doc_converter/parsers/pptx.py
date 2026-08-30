# Copyright (c) 2026 Chrys. All rights reserved.

"""PPTX parser — converts PowerPoint presentations to Markdown using python-pptx."""

from __future__ import annotations

from typing import Any

from chrys.service.tools.builtins.doc_converter.parsers._table import rows_to_markdown_table
from chrys.service.tools.builtins.doc_converter.parsers.base import (
    DocumentImageSink,
    ParsedDocument,
    VisualOccurrence,
)


class PptxParser:
    """Convert PPTX files to Markdown via ``python-pptx``."""

    @property
    def supported_extensions(self) -> frozenset[str]:
        return frozenset({".pptx"})

    def parse(self, path: str, *, image_sink: DocumentImageSink | None = None) -> ParsedDocument:
        from pptx import Presentation

        prs = Presentation(path)
        parts: list[str] = []
        picture_candidates: list[tuple[int, Any]] = []
        for i, slide in enumerate(prs.slides, 1):
            title = self._slide_title(slide)
            heading = f"# Slide {i}: {title}" if title else f"# Slide {i}"
            body_lines: list[str] = []
            for shape in slide.shapes:
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        cells = tuple(cell.text.strip() for cell in row.cells)
                        rows.append(cells)
                    if rows:
                        body_lines.append(rows_to_markdown_table(rows))
                elif shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and text != title:
                            body_lines.append(text)
            body = "\n\n".join(body_lines) if body_lines else "(no text content)"
            parts.append(f"{heading}\n\n{body}")
            if image_sink is not None:
                from pptx.shapes.group import GroupShape
                from pptx.shapes.picture import Picture

                def collect_pictures(shapes: Any, slide_number: int) -> None:
                    for candidate in shapes:
                        if isinstance(candidate, Picture):
                            picture_candidates.append((slide_number, candidate))
                        elif isinstance(candidate, GroupShape):
                            collect_pictures(candidate.shapes, slide_number)

                collect_pictures(slide.shapes, i)

        visuals: list[VisualOccurrence] = []
        slide_ordinals: dict[int, int] = {}
        if image_sink is not None:
            for candidate_index, (slide_number, shape) in enumerate(picture_candidates):
                if not image_sink.try_reserve_occurrence():
                    image_sink.record_unprocessed_occurrences(len(picture_candidates) - candidate_index)
                    break
                ordinal = slide_ordinals.get(slide_number, 0) + 1
                slide_ordinals[slide_number] = ordinal
                try:
                    image = shape.image
                    blob = image.blob
                except ValueError as exc:
                    if str(exc) == "no embedded image":
                        image_sink.record_linked_image()
                    else:
                        image_sink.record_candidate_failure()
                    continue
                except Exception:
                    image_sink.record_candidate_failure()
                    continue
                occurrence = image_sink.save_image(
                    blob,
                    location=f"Slide {slide_number}",
                    ordinal=ordinal,
                    source_name=image.filename,
                )
                if occurrence is not None:
                    visuals.append(occurrence)

        return ParsedDocument(
            markdown="\n\n".join(parts),
            visuals=tuple(visuals),
            warnings=image_sink.warnings if image_sink is not None else (),
        )

    @staticmethod
    def _slide_title(slide) -> str:
        """Extract the slide title, if any."""
        if slide.shapes.title is not None:
            return slide.shapes.title.text.strip()
        return ""
