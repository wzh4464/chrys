# Copyright (c) 2026 Chrys. All rights reserved.

"""Tests for doc_converter tool — convert_document + individual parsers."""

from __future__ import annotations

import asyncio
import hashlib
import json
import os
import threading
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from chrys.foundation.config.settings import SESSION_ROOT_DIR_ENV_VAR
from chrys.foundation.platform import get_platform
from chrys.foundation.text.images import MAX_IMAGE_BYTES as MAX_MODEL_IMAGE_BYTES
from chrys.foundation.text.images import MAX_IMAGE_SOURCE_BYTES
from chrys.foundation.tool_kinds import get_tool_kind
from chrys.foundation.tool_result_metadata import (
    TOOL_ERROR_DETAILS_METADATA_KEY,
    TOOL_ERROR_KIND_METADATA_KEY,
    TOOL_FAILED_METADATA_KEY,
)
from chrys.foundation.util.session_ids import session_short_id
from chrys.kernel import Message
from chrys.kernel.loop import LoopRecorder
from chrys.kernel.tools import SyncToolCancelledAfterCompletion
from chrys.service.state.serializers import serialized_message_payload
from chrys.service.state.store import JsonFileStateStore
from chrys.service.tools.builtins.doc_converter import (
    _TOKEN_THRESHOLD,
    DocConverterTools,
    _extract_toc,
    _normalize_document_text,
    _normalize_parsed_document,
    _render_visual_entries,
    _write_file,
    _write_unique_markdown,
)
from chrys.service.tools.builtins.doc_converter.artifacts import (
    MAX_ARTIFACT_BASENAME_BYTES,
    MAX_IMAGE_OCCURRENCES,
    MAX_SESSION_IMAGE_BYTES,
    MAX_SESSION_IMAGE_FILES,
    MAX_SOURCE_STEM_BYTES,
    MAX_STORED_IMAGE_BYTES,
    MAX_TOTAL_IMAGE_BYTES,
    MAX_UNIQUE_IMAGES,
    DocumentImageSink,
)
from chrys.service.tools.builtins.doc_converter.parsers.base import ParsedDocument, VisualOccurrence
from chrys.service.tools.builtins.filesystem import FilesystemTools, read_file
from chrys.service.tools.result_metadata import tool_result_metadata
from chrys.service.tools.session_artifacts import (
    resolve_document_image_artifact_handle,
    resolve_document_markdown_artifact_handle,
)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_PATCH_REGISTRY = "chrys.service.tools.builtins.doc_converter.registry"
_PATCH_TOK = "chrys.service.tools.builtins.doc_converter._tokenizer"
_PATCH_TOOL = "chrys.service.tools.builtins.doc_converter"


@dataclass
class _FakePlatformInfo:
    config_dir: Path
    os_name: str = "macos"
    arch: str = "arm64"
    shell: object = None
    data_dir: Path | None = None
    extra_shells: tuple = ()


@dataclass
class _FakeRuntime:
    platform: _FakePlatformInfo
    cwd: str = "/tmp"
    working_dirs: list = None

    def __post_init__(self):
        if self.working_dirs is None:
            self.working_dirs = []


def _make_runtime(tmp_path: Path) -> _FakeRuntime:
    """Create a minimal fake SessionEnvironment-like object."""
    platform = _FakePlatformInfo(config_dir=tmp_path)
    return _FakeRuntime(platform=platform)


class _FakeParser:
    """Fake parser for testing convert_document dispatch."""

    def __init__(
        self, text_content: str = "", *, side_effect: Exception | None = None, extensions: frozenset[str] | None = None
    ):
        self._text = text_content
        self._side_effect = side_effect
        self._extensions = extensions or frozenset({".pdf"})

    @property
    def supported_extensions(self) -> frozenset[str]:
        return self._extensions

    def parse(self, path: str, *, image_sink: DocumentImageSink | None = None) -> ParsedDocument:
        if self._side_effect is not None:
            raise self._side_effect
        return ParsedDocument(markdown=self._text)


class _VisualFakeParser(_FakeParser):
    """Fake parser that emits one image only when a sink is supplied."""

    def __init__(self, image_data: bytes, text_content: str = "# Report\n\ncontent") -> None:
        super().__init__(text_content)
        self._image_data = image_data
        self.last_sink: DocumentImageSink | None = None

    def parse(self, path: str, *, image_sink: DocumentImageSink | None = None) -> ParsedDocument:
        self.last_sink = image_sink
        if image_sink is None:
            return ParsedDocument(markdown=self._text)
        assert image_sink.try_reserve_occurrence()
        occurrence = image_sink.save_image(
            self._image_data,
            location="Page 1",
            ordinal=1,
            source_name="embedded.png",
        )
        visuals = (occurrence,) if occurrence is not None else ()
        return ParsedDocument(markdown=self._text, visuals=visuals, warnings=image_sink.warnings)


def _png_bytes(color: tuple[int, int, int] = (255, 0, 0)) -> bytes:
    from PIL import Image

    output = BytesIO()
    Image.new("RGB", (8, 8), color).save(output, format="PNG")
    return output.getvalue()


def _visual_paths(result: str) -> list[str]:
    paths: list[str] = []
    for line in result.splitlines():
        marker = line.find('{"path"')
        if marker >= 0:
            payload = json.loads(line[marker:])
            paths.append(payload["path"])
    return paths


def _session_artifact_path(session_dir: Path, reference: str) -> Path:
    resolved = resolve_document_image_artifact_handle(reference, session_dir)
    assert resolved is not None
    return Path(resolved)


def _session_markdown_artifact_path(session_dir: Path, reference: str) -> Path:
    resolved = resolve_document_markdown_artifact_handle(reference, session_dir)
    assert resolved is not None
    return Path(resolved)


def _saved_markdown_path(result: str, session_dir: Path) -> Path:
    handle = result.split("Saved Markdown handle: ", 1)[1].splitlines()[0]
    return _session_markdown_artifact_path(session_dir, handle)


def _saved_markdown_absolute_path(result: str) -> Path:
    path = result.split("Saved Markdown path: ", 1)[1].splitlines()[0]
    return Path(path)


def test_convert_document_contract_limits_visual_assets_to_pdf_docx_and_pptx() -> None:
    description = DocConverterTools.convert_document.description

    assert "embedded raster assets from PDF, DOCX, and PPTX" in description
    assert "XLS/XLSX conversion remains text/table-only" in description


def test_document_image_storage_limits_align_with_session_and_model_boundaries() -> None:
    assert MAX_UNIQUE_IMAGES == MAX_IMAGE_OCCURRENCES == 128
    assert MAX_SESSION_IMAGE_FILES == 8192
    assert MAX_STORED_IMAGE_BYTES == MAX_IMAGE_SOURCE_BYTES == 50 * 1024 * 1024
    assert MAX_TOTAL_IMAGE_BYTES == MAX_SESSION_IMAGE_BYTES == 512 * 1024 * 1024
    assert MAX_MODEL_IMAGE_BYTES == 3 * 1024 * 1024


def test_normalize_document_text_repairs_pairs_and_replaces_lone_surrogates() -> None:
    cases = {
        "ordinary text 🌍": "ordinary text 🌍",
        "split pair \ud83c\udf0d": "split pair 🌍",
        "lone high \ud83c": "lone high �",
        "lone low \udf0d": "lone low �",
        r"literal escape \ud83c": r"literal escape \ud83c",
    }

    for source, expected in cases.items():
        normalized = _normalize_document_text(source)
        normalized.encode("utf-8")
        assert normalized == expected


def test_normalize_parsed_document_covers_every_string_field() -> None:
    parsed = ParsedDocument(
        markdown="body \ud83c",
        visuals=(
            VisualOccurrence(
                location="Page \ud83c",
                ordinal=1,
                reference="chrys-session-document:image-\ud83c.png",
                source_name="source-\ud83c.png",
            ),
        ),
        warnings=("warning \ud83c",),
    )

    normalized = _normalize_parsed_document(parsed)

    normalized.markdown.encode("utf-8")
    normalized.visuals[0].location.encode("utf-8")
    normalized.visuals[0].reference.encode("utf-8")
    normalized.visuals[0].source_name.encode("utf-8")
    normalized.warnings[0].encode("utf-8")
    assert "�" in normalized.markdown


def test_visual_path_arguments_are_valid_json_for_windows_paths() -> None:
    reference = r'C:\Users\name\a "quoted" image.png'

    entry = _render_visual_entries((VisualOccurrence(location="Slide 1", ordinal=1, reference=reference),))[0]
    payload = json.loads(entry[entry.index("{") :])

    assert payload == {"path": reference}


def test_write_file_is_total_against_unpaired_surrogates(tmp_path: Path) -> None:
    target = tmp_path / "converted.md"

    _write_file(str(target), "damaged \ud83c text")

    assert target.read_text(encoding="utf-8") == r"damaged \ud83c text"


def test_document_image_sink_deduplicates_and_commits_returned_occurrences(tmp_path: Path) -> None:
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="report")
    image = _png_bytes()
    assert sink.try_reserve_occurrence()
    first = sink.save_image(image, location="Page 1", ordinal=1, source_name="first.png")
    assert sink.try_reserve_occurrence()
    second = sink.save_image(image, location="Page 2", ordinal=1, source_name="second.png")

    assert first is not None
    assert second is not None
    assert first.reference == second.reference
    files = list((tmp_path / "session" / "doc_converter").glob("*.png"))
    assert len(files) == 1

    sink.commit_occurrences((first, second))

    assert files[0].exists()


def test_document_image_sink_reuses_committed_content_and_cleans_new_abort(tmp_path: Path) -> None:
    root = tmp_path / "session" / "doc_converter"
    first_sink = DocumentImageSink(root, source_stem="report")
    assert first_sink.try_reserve_occurrence()
    existing = first_sink.save_image(_png_bytes(), location="Page 1", ordinal=1, source_name="existing.png")
    assert existing is not None
    first_sink.commit_occurrences((existing,))

    second_sink = DocumentImageSink(root, source_stem="report")
    assert second_sink.try_reserve_occurrence()
    second_copy = second_sink.save_image(_png_bytes(), location="Page 1", ordinal=1, source_name="copy.png")
    assert second_sink.try_reserve_occurrence()
    created = second_sink.save_image(_png_bytes((0, 0, 255)), location="Page 2", ordinal=1, source_name="new.png")
    assert second_copy is not None
    assert created is not None
    assert second_copy.reference == existing.reference

    second_sink.abort()

    assert _session_artifact_path(root.parent, existing.reference).exists()
    assert not _session_artifact_path(root.parent, created.reference).exists()


def test_aborting_creator_cannot_delete_image_committed_by_concurrent_conversion(tmp_path: Path) -> None:
    root = tmp_path / "session" / "doc_converter"
    creator = DocumentImageSink(root, source_stem="report")
    consumer = DocumentImageSink(root, source_stem="report")
    image = _png_bytes()
    assert creator.try_reserve_occurrence()
    created = creator.save_image(image, location="Page 1", ordinal=1, source_name="logo.png")
    assert consumer.try_reserve_occurrence()
    committed = consumer.save_image(image, location="Page 1", ordinal=1, source_name="logo.png")
    assert created is not None
    assert committed is not None
    assert created.reference == committed.reference

    consumer.commit_occurrences((committed,))
    creator.abort()

    assert _session_artifact_path(root.parent, committed.reference).exists()


def test_content_addressed_path_mismatch_is_reported_as_persistence_failure(tmp_path: Path) -> None:
    root = tmp_path / "session" / "doc_converter"
    creator = DocumentImageSink(root, source_stem="report")
    assert creator.try_reserve_occurrence()
    committed = creator.save_image(_png_bytes(), location="Page 1", ordinal=1, source_name="logo.png")
    assert committed is not None
    creator.commit_occurrences((committed,))
    artifact = _session_artifact_path(root.parent, committed.reference)
    artifact.write_bytes(_png_bytes((0, 0, 255)))

    collider = DocumentImageSink(root, source_stem="report")
    assert collider.try_reserve_occurrence()
    rejected = collider.save_image(_png_bytes(), location="Page 1", ordinal=1, source_name="logo.png")

    assert rejected is None
    assert collider.warnings == (
        "Skipped 1 image candidate(s) that could not be persisted or referenced as session artifacts.",
    )
    assert "decoded or normalized" not in collider.warnings[0]
    assert _session_artifact_path(root.parent, committed.reference).exists()


def test_image_handle_failure_is_reported_as_persistence_failure_and_cleans_orphan(tmp_path: Path) -> None:
    root = tmp_path / "session" / "doc_converter"
    sink = DocumentImageSink(root, source_stem="report")
    assert sink.try_reserve_occurrence()

    with patch.object(DocumentImageSink, "_model_reference", side_effect=ValueError("invalid handle")):
        occurrence = sink.save_image(_png_bytes(), location="Page 1", ordinal=1, source_name="logo.png")

    assert occurrence is None
    assert sink.warnings == (
        "Skipped 1 image candidate(s) that could not be persisted or referenced as session artifacts.",
    )
    assert len(list(root.glob("*.png"))) == 1

    sink.commit_occurrences(())

    assert list(root.glob("*.png")) == []


def test_document_image_sink_unique_limit_still_allows_repeated_digest(tmp_path: Path) -> None:
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="report")
    first: VisualOccurrence | None = None
    with patch("chrys.service.tools.builtins.doc_converter.artifacts.MAX_UNIQUE_IMAGES", 2):
        for index in range(2):
            assert sink.try_reserve_occurrence()
            occurrence = sink.save_image(
                _png_bytes((index, 0, 0)),
                location="Page 1",
                ordinal=index + 1,
                source_name=f"{index}.png",
            )
            assert occurrence is not None
            first = first or occurrence

        assert sink.try_reserve_occurrence()
        rejected = sink.save_image(
            _png_bytes((2, 0, 0)),
            location="Page 1",
            ordinal=3,
            source_name="over.png",
        )
        assert rejected is None
        assert first is not None
        assert sink.try_reserve_occurrence()
        repeated = sink.save_image(
            _png_bytes((0, 0, 0)),
            location="Page 2",
            ordinal=1,
            source_name="repeat.png",
        )

    assert repeated is not None
    assert repeated.reference == first.reference
    assert len(list((tmp_path / "session" / "doc_converter").iterdir())) == 2


def test_document_image_sink_enforces_occurrence_limit_with_bounded_warning(tmp_path: Path) -> None:
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="report")

    assert all(sink.try_reserve_occurrence() for _ in range(MAX_IMAGE_OCCURRENCES))
    assert sink.try_reserve_occurrence() is False
    sink.record_unprocessed_occurrences(72)

    assert len(sink.warnings) == 1
    assert "72 image candidate(s)" in sink.warnings[0]


def test_document_image_sink_enforces_total_normalized_byte_limit(tmp_path: Path) -> None:
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="report")

    with (
        patch(
            "chrys.service.tools.builtins.doc_converter.artifacts.MAX_TOTAL_IMAGE_BYTES",
            10,
        ),
        patch.object(
            DocumentImageSink,
            "_normalize_image",
            side_effect=[(b"first!", ".png"), (b"second", ".png")],
        ),
    ):
        assert sink.try_reserve_occurrence()
        first = sink.save_image(b"source-1", location="Page 1", ordinal=1, source_name="one.png")
        assert sink.try_reserve_occurrence()
        second = sink.save_image(b"source-2", location="Page 1", ordinal=2, source_name="two.png")

    assert first is not None
    assert second is None
    assert any("storage limits" in warning for warning in sink.warnings)


def test_document_image_sink_enforces_session_quota_and_still_reuses_digest(tmp_path: Path) -> None:
    root = tmp_path / "session" / "doc_converter"
    first_bytes = _png_bytes()
    second_bytes = _png_bytes((0, 0, 255))

    with patch(
        "chrys.service.tools.builtins.doc_converter.artifacts.MAX_SESSION_IMAGE_BYTES",
        len(first_bytes),
    ):
        first_sink = DocumentImageSink(root, source_stem="first")
        assert first_sink.try_reserve_occurrence()
        first = first_sink.save_image(first_bytes, location="Page 1", ordinal=1, source_name="first.png")
        assert first is not None
        first_sink.commit_occurrences((first,))

        second_sink = DocumentImageSink(root, source_stem="second")
        assert second_sink.try_reserve_occurrence()
        rejected = second_sink.save_image(
            second_bytes,
            location="Page 1",
            ordinal=1,
            source_name="second.png",
        )
        assert rejected is None
        assert any("session image storage limits" in warning for warning in second_sink.warnings)

        repeated_sink = DocumentImageSink(root, source_stem="repeat")
        assert repeated_sink.try_reserve_occurrence()
        repeated = repeated_sink.save_image(
            first_bytes,
            location="Page 1",
            ordinal=1,
            source_name="repeat.png",
        )

    assert repeated is not None
    assert repeated.reference == first.reference
    assert len(list(root.glob("*.png"))) == 1


def test_document_image_sink_enforces_session_file_count_with_streaming_scan(tmp_path: Path) -> None:
    root = tmp_path / "session" / "doc_converter"
    root.mkdir(parents=True)
    (root / "image-existing-one.png").write_bytes(b"one")
    (root / "image-existing-two.jpg").write_bytes(b"two")
    (root / "report.md").write_text("ignored", encoding="utf-8")
    sink = DocumentImageSink(root, source_stem="report")

    with (
        patch(
            "chrys.service.tools.builtins.doc_converter.artifacts.MAX_SESSION_IMAGE_FILES",
            2,
        ),
        patch.object(Path, "iterdir", side_effect=AssertionError("quota scan must stream")),
    ):
        assert sink.try_reserve_occurrence()
        occurrence = sink.save_image(
            _png_bytes(),
            location="Page 1",
            ordinal=1,
            source_name="new.png",
        )

    assert occurrence is None
    assert any("session image storage limits" in warning for warning in sink.warnings)
    assert sorted(path.name for path in root.glob("*.png")) == ["image-existing-one.png"]


def test_concurrent_image_sinks_cannot_oversubscribe_session_quota(tmp_path: Path) -> None:
    root = tmp_path / "session" / "doc_converter"
    payloads = (_png_bytes(), _png_bytes((0, 0, 255)))
    barrier = threading.Barrier(2)
    results: list[tuple[VisualOccurrence | None, tuple[str, ...]]] = []
    results_lock = threading.Lock()

    def save(payload: bytes) -> None:
        sink = DocumentImageSink(root, source_stem="report")
        assert sink.try_reserve_occurrence()
        barrier.wait()
        occurrence = sink.save_image(payload, location="Page 1", ordinal=1, source_name="image.png")
        if occurrence is not None:
            sink.commit_occurrences((occurrence,))
        with results_lock:
            results.append((occurrence, sink.warnings))

    with patch(
        "chrys.service.tools.builtins.doc_converter.artifacts.MAX_SESSION_IMAGE_BYTES",
        max(map(len, payloads)),
    ):
        threads = [threading.Thread(target=save, args=(payload,)) for payload in payloads]
        for thread in threads:
            thread.start()
        for thread in threads:
            thread.join()

    assert len(results) == 2
    assert sum(occurrence is not None for occurrence, _warnings in results) == 1
    assert any("session image storage limits" in warning for _occurrence, warnings in results for warning in warnings)
    assert len(list(root.glob("*.png"))) == 1


def test_document_image_sink_skips_normalized_image_over_per_file_limit(tmp_path: Path) -> None:
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="report")

    with (
        patch("chrys.service.tools.builtins.doc_converter.artifacts.MAX_STORED_IMAGE_BYTES", 8),
        patch(
            "chrys.service.tools.builtins.doc_converter.artifacts.compress_image_data",
            return_value=b"x" * 9,
        ),
    ):
        assert sink.try_reserve_occurrence()
        occurrence = sink.save_image(b"unsupported-raster", location="Page 1", ordinal=1, source_name="huge.bmp")

    assert occurrence is None
    assert len(sink.warnings) == 1
    assert "could not be decoded or normalized" in sink.warnings[0]
    assert not (tmp_path / "session" / "doc_converter").exists()


def test_document_image_and_markdown_names_are_utf8_byte_bounded(tmp_path: Path) -> None:
    long_stem = "报告" * 100
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem=long_stem)
    assert sink.try_reserve_occurrence()
    occurrence = sink.save_image(_png_bytes(), location="Page 1", ordinal=1, source_name="内部名.png")

    assert occurrence is not None
    image_name = _session_artifact_path(tmp_path / "session", occurrence.reference).name
    assert len(image_name.encode("utf-8")) <= MAX_ARTIFACT_BASENAME_BYTES
    assert image_name == f"image-{hashlib.sha256(_png_bytes()).hexdigest()}.png"
    assert long_stem not in image_name
    assert "内部名" not in image_name


@pytest.mark.skipif(get_platform().is_windows, reason="symlink creation is not generally available on Windows")
def test_document_image_sink_rejects_redirected_artifact_root(tmp_path: Path) -> None:
    session_dir = tmp_path / "session"
    outside = tmp_path / "outside"
    session_dir.mkdir()
    outside.mkdir()
    (session_dir / "doc_converter").symlink_to(outside, target_is_directory=True)
    sink = DocumentImageSink(session_dir / "doc_converter", source_stem="report")
    assert sink.try_reserve_occurrence()

    occurrence = sink.save_image(_png_bytes(), location="Page 1", ordinal=1, source_name="pixel.png")

    assert occurrence is None
    assert list(outside.iterdir()) == []
    assert any("artifact directory" in warning for warning in sink.warnings)


@pytest.mark.skipif(get_platform().is_windows, reason="symlink creation is not generally available on Windows")
def test_markdown_writer_canonicalizes_trusted_session_parent_symlink(tmp_path: Path) -> None:
    real_root = tmp_path / "real"
    real_root.mkdir()
    alias_root = tmp_path / "alias"
    alias_root.symlink_to(real_root, target_is_directory=True)
    aliased_session_dir = alias_root / "session"

    saved = _write_unique_markdown(
        os.fspath(aliased_session_dir / "doc_converter"),
        "report",
        "# Report",
    )

    expected = real_root / "session" / "doc_converter" / "report.md"
    assert saved.handle == "chrys-session-document:report.md"
    assert saved.path == os.fspath(expected)
    assert expected.read_text(encoding="utf-8") == "# Report"
    assert resolve_document_markdown_artifact_handle(saved.handle, aliased_session_dir) == os.fspath(expected)


@pytest.mark.skipif(get_platform().is_windows, reason="symlink creation is not generally available on Windows")
def test_markdown_writer_rejects_redirected_artifact_child(tmp_path: Path) -> None:
    session_dir = tmp_path / "session"
    outside = tmp_path / "outside"
    session_dir.mkdir()
    outside.mkdir()
    (session_dir / "doc_converter").symlink_to(outside, target_is_directory=True)

    with pytest.raises(OSError, match="must not be redirected"):
        _write_unique_markdown(os.fspath(session_dir / "doc_converter"), "report", "# Report")

    assert list(outside.iterdir()) == []


@pytest.mark.skipif(get_platform().is_windows, reason="symlink creation is not generally available on Windows")
def test_document_image_sink_rejects_redirected_existing_image(tmp_path: Path) -> None:
    root = tmp_path / "session" / "doc_converter"
    first_sink = DocumentImageSink(root, source_stem="report")
    assert first_sink.try_reserve_occurrence()
    first = first_sink.save_image(_png_bytes(), location="Page 1", ordinal=1, source_name="pixel.png")
    assert first is not None
    first_sink.abort()

    outside = tmp_path / "outside.png"
    outside.write_bytes(b"outside")
    artifact = _session_artifact_path(tmp_path / "session", first.reference)
    artifact.symlink_to(outside)
    second_sink = DocumentImageSink(root, source_stem="report")
    assert second_sink.try_reserve_occurrence()

    occurrence = second_sink.save_image(_png_bytes(), location="Page 1", ordinal=1, source_name="pixel.png")

    assert occurrence is None
    assert outside.read_bytes() == b"outside"
    assert second_sink.warnings == (
        "Skipped 1 image candidate(s) that could not be persisted or referenced as session artifacts.",
    )


def test_concurrent_same_name_image_writes_do_not_overwrite_different_bytes(tmp_path: Path) -> None:
    root = tmp_path / "session" / "doc_converter"
    barrier = threading.Barrier(2)
    results: list[VisualOccurrence] = []
    result_lock = threading.Lock()

    def save(color: tuple[int, int, int]) -> None:
        sink = DocumentImageSink(root, source_stem="report")
        assert sink.try_reserve_occurrence()
        barrier.wait()
        occurrence = sink.save_image(_png_bytes(color), location="Page 1", ordinal=1, source_name="same.png")
        assert occurrence is not None
        sink.commit_occurrences((occurrence,))
        with result_lock:
            results.append(occurrence)

    threads = [
        threading.Thread(target=save, args=((255, 0, 0),)),
        threading.Thread(target=save, args=((0, 0, 255),)),
    ]
    for thread in threads:
        thread.start()
    for thread in threads:
        thread.join()

    assert len(results) == 2
    assert results[0].reference != results[1].reference
    stored = sorted(root.glob("*.png"))
    assert len(stored) == 2
    assert {path.read_bytes() for path in stored} == {_png_bytes((255, 0, 0)), _png_bytes((0, 0, 255))}


# ---------------------------------------------------------------------------
# TOC extraction
# ---------------------------------------------------------------------------


def test_extract_toc_basic() -> None:
    md = "# Title\n\nSome text.\n\n## Section A\n\nMore text.\n\n### Subsection\n\n## Section B\n"
    toc = _extract_toc(md)
    assert "1|- Title" in toc
    assert "5|  - Section A" in toc
    assert "9|    - Subsection" in toc
    assert "11|  - Section B" in toc


def test_extract_toc_no_headings() -> None:
    toc = _extract_toc("Just plain text.\nNo headings here.\n")
    assert toc == "(no headings found)"


def test_extract_toc_skips_empty_hashes() -> None:
    toc = _extract_toc("#\n## Real heading\n###\n")
    assert "Real heading" in toc
    # Bare "#" or "###" with no title text should not appear
    assert toc.count("- ") == 1


def test_extract_toc_truncates_when_large() -> None:
    """TOC with many headings gets truncated at _TOC_MAX_TOKENS budget."""
    lines = [f"# Heading number {i}" for i in range(500)]
    md = "\n\nSome text.\n\n".join(lines)

    with patch(_PATCH_TOK) as mock_tok:
        mock_tok.count_tokens.return_value = 50  # 50 tokens per entry -> ~100 fit in 5k
        toc = _extract_toc(md)

    assert "truncated" in toc.lower()
    assert "heading" in toc.lower()
    assert "Heading number 0" in toc
    assert "Heading number 499" not in toc


# ---------------------------------------------------------------------------
# convert_document — small document (inline return)
# ---------------------------------------------------------------------------


async def test_convert_small_document(tmp_path: Path) -> None:
    """Small doc (< threshold) returns full markdown inline."""
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake-content")

    small_md = "# Planet \ud83c\udf0d\n\nDamaged \ud83c text.\n"
    fake_parser = _FakeParser(small_md)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser):
        runtime = _make_runtime(tmp_path)
        tools = DocConverterTools(runtime, session_id="test123")
        result = await tools.convert_document(str(doc))

    result.encode("utf-8")
    assert "report.pdf" in result
    assert "# Planet 🌍" in result
    assert "Damaged � text." in result
    assert "lines" in result
    assert "tokens" in result

    async def on_checkpoint() -> None:
        return None

    recorder = LoopRecorder(on_checkpoint=on_checkpoint, message_hasher=serialized_message_payload)
    await recorder.record_pre_call([Message(role="tool", contents=[result])])


async def test_convert_small_document_escapes_surrogate_path_but_parses_raw_path(tmp_path: Path) -> None:
    resolved = "/work/report-\udcff.pdf"
    fake_parser = MagicMock()
    fake_parser.parse.return_value = ParsedDocument(markdown="# Report\n\ncontent\n")
    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime, session_id="test123")

    with (
        patch(f"{_PATCH_TOOL}.resolve_existing_path", return_value=resolved),
        patch(f"{_PATCH_TOOL}.os.path.isfile", return_value=True),
        patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser),
    ):
        result = await tools.convert_document("report.pdf")

    result.encode("utf-8")
    assert r"File: /work/report-\udcff.pdf" in result
    fake_parser.parse.assert_called_once_with(resolved, image_sink=None)


async def test_convert_document_falls_back_to_narrow_no_break_space_filename(tmp_path: Path) -> None:
    """convert_document tolerates macOS screenshot U+202F whitespace in filenames."""
    real = tmp_path / "report 9.00.00\u202fAM.pdf"
    real.write_bytes(b"%PDF-fake-content")
    requested = str(tmp_path / "report 9.00.00 AM.pdf")

    small_md = "# Report\n\ncontent\n"
    fake_parser = _FakeParser(small_md)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser):
        runtime = _make_runtime(tmp_path)
        tools = DocConverterTools(runtime, session_id="test123")
        result = await tools.convert_document(requested)

    assert "# Report" in result
    assert str(real) in result


async def test_convert_document_default_gate_avoids_all_image_work(tmp_path: Path) -> None:
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    parser = _VisualFakeParser(_png_bytes())
    tools = DocConverterTools(_make_runtime(tmp_path), session_dir=tmp_path / "session")

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=parser):
        result = await tools.convert_document(str(doc))

    assert parser.last_sink is None
    assert "Extracted visual assets" not in result
    assert "view_image" not in result
    assert not (tmp_path / "session" / "doc_converter").exists()


async def test_convert_document_enabled_flag_without_session_still_avoids_image_work(tmp_path: Path) -> None:
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    parser = _VisualFakeParser(_png_bytes())
    tools = DocConverterTools(_make_runtime(tmp_path))
    tools.set_image_extraction_enabled(True)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=parser):
        result = await tools.convert_document(str(doc))

    assert parser.last_sink is None
    assert "Extracted visual assets" not in result
    assert "view_image" not in result


async def test_convert_document_disabled_image_only_pdf_keeps_page_placeholder(tmp_path: Path) -> None:
    from PIL import Image

    from chrys.service.tools.builtins.doc_converter.parsers.pdf import PdfParser

    doc = tmp_path / "image-only.pdf"
    Image.new("RGB", (8, 8), (255, 0, 0)).save(doc, format="PDF")
    session_dir = tmp_path / "session"
    tools = DocConverterTools(_make_runtime(tmp_path), session_dir=session_dir)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=PdfParser()):
        result = await tools.convert_document(str(doc))

    assert "# Page 1\n\n(no text content)" in result
    assert "view_image" not in result
    assert not (session_dir / "doc_converter").exists()


async def test_convert_document_enabled_gate_returns_viewable_image_path(tmp_path: Path) -> None:
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    parser = _VisualFakeParser(_png_bytes())
    session_dir = tmp_path / "session"
    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime, session_dir=session_dir)
    tools.set_image_extraction_enabled(True)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=parser):
        result = await tools.convert_document(str(doc))

    assert parser.last_sink is not None
    assert result.count("Use view_image") == 1
    paths = _visual_paths(result)
    assert len(paths) == 1
    assert paths[0].startswith("chrys-session-document:")
    image_result = FilesystemTools(runtime, session_dir=session_dir).view_image(paths[0])
    assert image_result[0].media_type == "image/png"


async def test_visual_index_participates_in_large_document_threshold(tmp_path: Path) -> None:
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    parser = _VisualFakeParser(_png_bytes(), text_content="# Short report")
    session_dir = tmp_path / "session"
    tools = DocConverterTools(_make_runtime(tmp_path), session_dir=session_dir)
    tools.set_image_extraction_enabled(True)

    def count_tokens(text: str) -> int:
        return _TOKEN_THRESHOLD if "Extracted visual assets" in text else 1

    with (
        patch(f"{_PATCH_REGISTRY}.get_parser", return_value=parser),
        patch(_PATCH_TOK) as mock_tok,
    ):
        mock_tok.count_tokens.side_effect = count_tokens
        result = await tools.convert_document(str(doc))

    assert "Document is too large to return inline" in result
    saved_path = _saved_markdown_path(result, session_dir)
    assert "Extracted visual assets" in saved_path.read_text(encoding="utf-8")


async def test_convert_document_empty_text_with_visual_does_not_drop_image(tmp_path: Path) -> None:
    doc = tmp_path / "empty.docx"
    doc.write_bytes(b"PK-fake")
    parser = _VisualFakeParser(_png_bytes(), text_content="")
    tools = DocConverterTools(_make_runtime(tmp_path), session_dir=tmp_path / "session")
    tools.set_image_extraction_enabled(True)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=parser):
        result = await tools.convert_document(str(doc))

    assert "document converted but produced no text content" not in result
    assert len(_visual_paths(result)) == 1


async def test_convert_document_failed_images_warn_without_view_image_instruction(tmp_path: Path) -> None:
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    parser = _VisualFakeParser(b"not an image")
    tools = DocConverterTools(_make_runtime(tmp_path), session_dir=tmp_path / "session")
    tools.set_image_extraction_enabled(True)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=parser):
        result = await tools.convert_document(str(doc))

    assert "Image extraction warnings" in result
    assert "view_image" not in result
    assert _visual_paths(result) == []


# ---------------------------------------------------------------------------
# convert_document — large document (save to disk)
# ---------------------------------------------------------------------------


async def test_convert_large_document(tmp_path: Path, monkeypatch) -> None:
    """Large doc (>= threshold) saves to session dir and returns TOC + metadata."""
    monkeypatch.delenv(SESSION_ROOT_DIR_ENV_VAR, raising=False)
    doc = tmp_path / "big.docx"
    doc.write_bytes(b"PK-fake-docx")

    lines = [f"# Chapter {i}\n\nParagraph content for chapter {i}. " + "word " * 200 for i in range(1, 51)]
    large_md = "\n".join(lines)

    fake_parser = _FakeParser(large_md, extensions=frozenset({".docx"}))

    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime, session_id="sess_abc")

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser), patch(_PATCH_TOK) as mock_tok:
        mock_tok.count_tokens.return_value = _TOKEN_THRESHOLD + 1000
        result = await tools.convert_document(str(doc))

    assert "too large to return inline" in result.lower() or "Saved Markdown handle:" in result
    assert "Table of Contents" in result
    assert "read_file" in result

    doc_converter_dir = tmp_path / "sessions" / session_short_id("sess_abc") / "doc_converter"
    saved_files = list(doc_converter_dir.glob("*.md"))
    assert len(saved_files) == 1
    assert saved_files[0].read_text(encoding="utf-8") == large_md


async def test_convert_large_document_prefers_supplied_session_dir(tmp_path: Path, monkeypatch) -> None:
    """A resolved session_dir should win over the environment fallback."""
    doc = tmp_path / "big.docx"
    doc.write_bytes(b"PK-fake-docx")
    large_md = "# Heading\n" + "content " * 5000
    fake_parser = _FakeParser(large_md, extensions=frozenset({".docx"}))
    explicit_session_dir = tmp_path / "explicit" / "sess_abc"
    monkeypatch.setenv(SESSION_ROOT_DIR_ENV_VAR, str(tmp_path / "env-root"))

    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime, session_id="sess_abc", session_dir=explicit_session_dir)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser), patch(_PATCH_TOK) as mock_tok:
        mock_tok.count_tokens.return_value = _TOKEN_THRESHOLD + 100
        await tools.convert_document(str(doc))

    saved_files = list((explicit_session_dir / "doc_converter").glob("*.md"))
    assert len(saved_files) == 1
    assert not (tmp_path / "env-root" / "sessions").exists()


async def test_convert_large_document_returns_path_usable_without_session_bound_read_file(tmp_path: Path) -> None:
    doc = tmp_path / "big.docx"
    doc.write_bytes(b"PK-fake-docx")
    markdown = "# Heading\n" + "content " * 5000
    parser = _FakeParser(markdown, extensions=frozenset({".docx"}))
    session_dir = tmp_path / "session"
    tools = DocConverterTools(_make_runtime(tmp_path), session_dir=session_dir)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=parser):
        result = await tools.convert_document(str(doc))

    saved_path = _saved_markdown_absolute_path(result)
    assert saved_path == session_dir / "doc_converter" / "big.md"
    assert "1|# Heading" in read_file(os.fspath(saved_path))
    assert "session-bound read_file" in result
    assert "filesystem or shell tool" in result


async def test_concurrent_same_name_large_conversions_keep_different_markdown(tmp_path: Path) -> None:
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    session_dir = tmp_path / "session"
    tools = DocConverterTools(_make_runtime(tmp_path), session_dir=session_dir)
    first = _FakeParser("# First\n\n" + "alpha " * 6000)
    second = _FakeParser("# Second\n\n" + "beta " * 6000)

    with patch(f"{_PATCH_REGISTRY}.get_parser", side_effect=[first, second]), patch(_PATCH_TOK) as mock_tok:
        mock_tok.count_tokens.return_value = _TOKEN_THRESHOLD + 1
        results = await asyncio.gather(
            tools.convert_document(str(doc)),
            tools.convert_document(str(doc)),
        )

    saved_paths = [_saved_markdown_path(result, session_dir) for result in results]
    assert saved_paths[0] != saved_paths[1]
    assert {path.read_text(encoding="utf-8") for path in saved_paths} == {first._text, second._text}


async def test_convert_large_document_uses_fixed_twenty_visual_preview(tmp_path: Path) -> None:
    class _ManyVisualParser(_FakeParser):
        def parse(self, path: str, *, image_sink: DocumentImageSink | None = None) -> ParsedDocument:
            assert image_sink is not None
            visuals: list[VisualOccurrence] = []
            for index in range(25):
                assert image_sink.try_reserve_occurrence()
                occurrence = image_sink.save_image(
                    _png_bytes((index, 0, 0)),
                    location=f"Page {index + 1}",
                    ordinal=1,
                    source_name=f"{index}.png",
                )
                assert occurrence is not None
                visuals.append(occurrence)
            return ParsedDocument(
                markdown="# Large report\n\ncontent",
                visuals=tuple(visuals),
                warnings=image_sink.warnings,
            )

    doc = tmp_path / "large.pdf"
    doc.write_bytes(b"%PDF-fake")
    session_dir = tmp_path / "session"
    tools = DocConverterTools(_make_runtime(tmp_path), session_dir=session_dir)
    tools.set_image_extraction_enabled(True)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=_ManyVisualParser()), patch(_PATCH_TOK) as mock_tok:
        mock_tok.count_tokens.return_value = _TOKEN_THRESHOLD + 1
        result = await tools.convert_document(str(doc))

    assert len(_visual_paths(result)) == 20
    assert "5 more image occurrence(s)" in result
    assert result.count("Use view_image") == 1
    saved_path = _saved_markdown_path(result, session_dir)
    saved = saved_path.read_text(encoding="utf-8")
    assert len(_visual_paths(saved)) == 25
    assert saved.count("Use view_image") == 1
    assert len(list((session_dir / "doc_converter").glob("*.png"))) == 25


async def test_convert_large_markdown_filename_is_utf8_byte_bounded(tmp_path: Path) -> None:
    resolved = "/work/" + ("报告" * 100) + ".pdf"
    parser = _FakeParser("# Report\n\n" + "content " * 6000)
    session_dir = tmp_path / "session"
    tools = DocConverterTools(_make_runtime(tmp_path), session_dir=session_dir)

    with (
        patch(f"{_PATCH_TOOL}.resolve_existing_path", return_value=resolved),
        patch(f"{_PATCH_TOOL}.os.path.isfile", return_value=True),
        patch(f"{_PATCH_REGISTRY}.get_parser", return_value=parser),
    ):
        result = await tools.convert_document("report.pdf")

    saved_path = _saved_markdown_path(result, session_dir)
    assert len(saved_path.name.encode("utf-8")) <= MAX_ARTIFACT_BASENAME_BYTES
    assert len(saved_path.stem.encode("utf-8")) <= MAX_SOURCE_STEM_BYTES
    assert saved_path.exists()


async def test_convert_large_document_normalizes_before_count_toc_and_save(tmp_path: Path) -> None:
    doc = tmp_path / "surrogate.pdf"
    doc.write_bytes(b"%PDF-fake")
    raw_markdown = "# Planet \ud83c\udf0d\n\nDamaged \ud83c text\n\n" + "content " * 6000
    fake_parser = _FakeParser(raw_markdown)
    runtime = _make_runtime(tmp_path)
    session_dir = tmp_path / "surrogate_session"
    tools = DocConverterTools(runtime, session_id="surrogate_session", session_dir=session_dir)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser):
        result = await tools.convert_document(str(doc))

    saved_files = list((session_dir / "doc_converter").glob("*.md"))
    assert len(saved_files) == 1
    saved = saved_files[0].read_text(encoding="utf-8")
    surrogate_positions = [index for index, char in enumerate(saved) if 0xD800 <= ord(char) <= 0xDFFF]
    assert not surrogate_positions, f"saved text contains surrogates at positions {surrogate_positions[:10]}"
    assert "# Planet 🌍" in saved
    assert "Damaged � text" in saved
    assert f"{len(saved)} chars" in result
    assert "Planet 🌍" in result
    result.encode("utf-8")


async def test_convert_large_document_reports_save_failure(tmp_path: Path) -> None:
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    fake_parser = _FakeParser("# Report\n\n" + "content " * 6000)
    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime, session_id="save_failure")
    metadata: dict[str, object] = {}
    token = tool_result_metadata.set(metadata)

    try:
        with (
            patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser),
            patch("chrys.service.tools.builtins.doc_converter._write_file", side_effect=OSError("disk full")),
        ):
            result = await tools.convert_document(str(doc))
    finally:
        tool_result_metadata.reset(token)

    assert result.startswith("Error:")
    assert "failed to save converted document" in result
    assert "disk full" in result
    assert metadata[TOOL_FAILED_METADATA_KEY] is True
    assert metadata[TOOL_ERROR_KIND_METADATA_KEY] == "document_save_failed"


async def test_convert_large_document_reports_handle_creation_failure(tmp_path: Path) -> None:
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    fake_parser = _FakeParser("# Report\n\n" + "content " * 6000)
    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime, session_id="handle_failure", session_dir=tmp_path / "session")

    with (
        patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser),
        patch(f"{_PATCH_TOOL}._unique_path", return_value="/session-\udcff/report.md"),
        patch(f"{_PATCH_TOOL}.make_document_artifact_handle", side_effect=ValueError("invalid artifact name")),
        patch(f"{_PATCH_TOOL}._write_file") as write_file_mock,
    ):
        result = await tools.convert_document(str(doc))

    result.encode("utf-8")
    assert result.startswith("Error:")
    assert "failed to save converted document" in result
    assert "invalid artifact name" in result
    write_file_mock.assert_not_called()


def test_document_converter_and_filesystem_reader_share_derived_session_dir(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.delenv(SESSION_ROOT_DIR_ENV_VAR, raising=False)
    import chrys.service.tools.session_artifacts as session_artifacts

    monkeypatch.setattr(
        session_artifacts,
        "resolve_sessions_dir",
        lambda config_dir, *, create: config_dir / "sessions",
    )
    runtime = _make_runtime(tmp_path / "config-\udcff")

    converter = DocConverterTools(runtime, session_id="shared_session")
    filesystem = FilesystemTools(runtime, session_id="shared_session")

    assert converter._session_dir == filesystem._session_dir
    assert converter._session_dir is not None
    assert str(converter._session_dir).startswith(str(runtime.platform.config_dir))


async def test_convert_large_document_creates_addressable_safe_output_path(tmp_path: Path) -> None:
    resolved = "/work/report-\udcff.pdf"
    markdown = "# Report\n\n" + "content " * 6000
    fake_parser = MagicMock()
    fake_parser.parse.return_value = ParsedDocument(markdown=markdown)
    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime, session_id="surrogate_path", session_dir=tmp_path / "session")

    with (
        patch(f"{_PATCH_TOOL}.resolve_existing_path", return_value=resolved),
        patch(f"{_PATCH_TOOL}.os.path.isfile", return_value=True),
        patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser),
    ):
        result = await tools.convert_document("report.pdf")

    result.encode("utf-8")
    assert r"File: /work/report-\udcff.pdf" in result
    returned_handle = result.split("Saved Markdown handle: ", 1)[1].splitlines()[0]
    expected_path = tmp_path / "session" / "doc_converter" / "report-_.md"
    assert _session_markdown_artifact_path(tmp_path / "session", returned_handle) == expected_path
    assert expected_path.read_text(encoding="utf-8") == markdown
    assert "1|# Report" in FilesystemTools(runtime, session_dir=tmp_path / "session").read_file(returned_handle)
    fake_parser.parse.assert_called_once_with(resolved, image_sink=None)


async def test_convert_large_document_returns_addressable_handle_for_surrogate_session_dir(
    tmp_path: Path,
) -> None:
    if not get_platform().is_linux:
        pytest.skip("invalid-UTF-8 filesystem paths are Linux-specific")

    raw_session_dir = os.path.join(os.fsencode(tmp_path), b"session-\xff")
    os.mkdir(raw_session_dir)
    session_dir = Path(os.fsdecode(raw_session_dir))
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    markdown = "# Report\n\n" + "content " * 6000
    fake_parser = _FakeParser(markdown)
    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime, session_id="surrogate_session_dir", session_dir=session_dir)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser):
        result = await tools.convert_document(str(doc))

    result.encode("utf-8")
    handle = result.split("Saved Markdown handle: ", 1)[1].splitlines()[0]
    expected_path = os.path.join(os.fsdecode(raw_session_dir), "doc_converter", "report.md")

    assert handle.isascii()
    assert "Saved Markdown path:" not in result
    read_result = FilesystemTools(runtime, session_dir=session_dir).read_file(handle)
    read_result.encode("utf-8")
    assert "1|# Report" in read_result
    assert os.fsencode(expected_path) == os.path.join(raw_session_dir, b"doc_converter", b"report.md")
    assert Path(expected_path).read_text(encoding="utf-8") == markdown


async def test_convert_document_returns_image_handle_for_surrogate_session_dir(tmp_path: Path) -> None:
    if not get_platform().is_linux:
        pytest.skip("invalid-UTF-8 filesystem paths are Linux-specific")

    raw_session_dir = os.path.join(os.fsencode(tmp_path), b"image-session-\xff")
    os.mkdir(raw_session_dir)
    session_dir = Path(os.fsdecode(raw_session_dir))
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    parser = _VisualFakeParser(_png_bytes())
    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime, session_dir=session_dir)
    tools.set_image_extraction_enabled(True)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=parser):
        result = await tools.convert_document(str(doc))

    result.encode("utf-8")
    paths = _visual_paths(result)
    assert len(paths) == 1
    assert paths[0].startswith("chrys-session-document:")
    image_result = FilesystemTools(runtime, session_dir=session_dir).view_image(paths[0])
    assert image_result[0].media_type == "image/png"
    source_path = image_result[0].additional_properties["source_path"]
    source_path.encode("utf-8")
    assert r"image-session-\udcff" in source_path

    message = Message(role="tool", contents=image_result)
    wire_payload = serialized_message_payload(message)
    wire_payload.encode("utf-8")
    store = JsonFileStateStore(tmp_path / "persisted-sessions")
    await store.save_session("surrogate-image-metadata", {"messages": [message], "compressed_msgs": []})
    restored = await store.load_session("surrogate-image-metadata")

    assert restored is not None
    restored_content = restored["messages"][0].contents[0]
    assert restored_content.additional_properties["source_path"] == source_path


# ---------------------------------------------------------------------------
# convert_document — no session_id (large doc, no save dir)
# ---------------------------------------------------------------------------


async def test_convert_large_no_session(tmp_path: Path) -> None:
    """Large doc with no session_id returns TOC without saving."""
    doc = tmp_path / "nosess.pdf"
    doc.write_bytes(b"%PDF-fake")

    large_md = "# Heading\n" + "content " * 5000
    fake_parser = _FakeParser(large_md)

    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime, session_id="")  # no session

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser), patch(_PATCH_TOK) as mock_tok:
        mock_tok.count_tokens.return_value = _TOKEN_THRESHOLD + 500
        result = await tools.convert_document(str(doc))

    assert "No session directory" in result
    assert "Table of Contents" in result


# ---------------------------------------------------------------------------
# Error cases
# ---------------------------------------------------------------------------


async def test_convert_file_not_found(tmp_path: Path) -> None:
    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime)
    result = await tools.convert_document(str(tmp_path / "nonexistent.pdf"))
    assert result.startswith("Error:")
    assert "file not found" in result


async def test_convert_file_not_found_escapes_display_path_but_keeps_raw_metadata(tmp_path: Path) -> None:
    resolved = "/work/missing-\udcff.pdf"
    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime)
    metadata: dict[str, object] = {}
    token = tool_result_metadata.set(metadata)

    try:
        with (
            patch(f"{_PATCH_TOOL}.resolve_existing_path", return_value=None),
            patch(f"{_PATCH_TOOL}.resolve_workspace_path", return_value=resolved),
            patch(f"{_PATCH_TOOL}.os.path.isdir", return_value=False),
        ):
            result = await tools.convert_document("missing.pdf")
    finally:
        tool_result_metadata.reset(token)

    result.encode("utf-8")
    assert r"file not found — /work/missing-\udcff.pdf" in result
    details = metadata[TOOL_ERROR_DETAILS_METADATA_KEY]
    assert isinstance(details, dict)
    assert details["resolved_path"] == resolved


async def test_convert_directory(tmp_path: Path) -> None:
    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime)
    result = await tools.convert_document(str(tmp_path))
    assert result.startswith("Error:")
    assert "directory" in result


async def test_convert_unsupported_format(tmp_path: Path) -> None:
    f = tmp_path / "image.png"
    f.write_bytes(b"\x89PNG")

    with (
        patch(f"{_PATCH_REGISTRY}.get_parser", return_value=None),
        patch(f"{_PATCH_REGISTRY}.supported_extensions", return_value=frozenset({".pdf"})),
    ):
        runtime = _make_runtime(tmp_path)
        tools = DocConverterTools(runtime)
        result = await tools.convert_document(str(f))

    assert result.startswith("Error:")
    assert "unsupported" in result.lower()


async def test_convert_import_error(tmp_path: Path) -> None:
    """When parser deps are not installed, returns a clear error."""
    doc = tmp_path / "doc.pdf"
    doc.write_bytes(b"%PDF-fake")

    fake_parser = _FakeParser(side_effect=ImportError("No module named 'pypdf'"))

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser):
        runtime = _make_runtime(tmp_path)
        tools = DocConverterTools(runtime)
        result = await tools.convert_document(str(doc))

    assert result.startswith("Error:")
    assert "missing dependencies" in result.lower()


async def test_convert_empty_result(tmp_path: Path) -> None:
    """Document that converts to empty text returns informative message."""
    doc = tmp_path / "empty.pdf"
    doc.write_bytes(b"%PDF-empty")

    fake_parser = _FakeParser("")

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser):
        runtime = _make_runtime(tmp_path)
        tools = DocConverterTools(runtime)
        result = await tools.convert_document(str(doc))

    assert "no text content" in result.lower()


async def test_convert_exception(tmp_path: Path) -> None:
    """Conversion exception returns error string."""
    doc = tmp_path / "bad.pdf"
    doc.write_bytes(b"%PDF-corrupt")

    fake_parser = _FakeParser(side_effect=RuntimeError("corrupt PDF"))

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser):
        runtime = _make_runtime(tmp_path)
        tools = DocConverterTools(runtime)
        result = await tools.convert_document(str(doc))

    assert result.startswith("Error:")
    assert "corrupt PDF" in result


async def test_cancelled_completed_conversion_carries_final_result_and_keeps_images(tmp_path: Path) -> None:
    started = threading.Event()
    release = threading.Event()
    image = _png_bytes()

    class _BlockingParser(_FakeParser):
        def parse(self, path: str, *, image_sink: DocumentImageSink | None = None) -> ParsedDocument:
            assert image_sink is not None
            assert image_sink.try_reserve_occurrence()
            occurrence = image_sink.save_image(
                image,
                location="Page 1",
                ordinal=1,
                source_name="pixel.png",
            )
            assert occurrence is not None
            started.set()
            assert release.wait(timeout=5)
            return ParsedDocument(markdown="# Report", visuals=(occurrence,), warnings=image_sink.warnings)

    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    session_dir = tmp_path / "session"
    runtime = _make_runtime(tmp_path)
    tools = DocConverterTools(runtime, session_dir=session_dir)
    tools.set_image_extraction_enabled(True)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=_BlockingParser()):
        task = asyncio.create_task(tools.convert_document(str(doc)))
        assert await asyncio.to_thread(started.wait, 5)
        task.cancel()
        release.set()
        with pytest.raises(SyncToolCancelledAfterCompletion) as exc_info:
            await task

    completed = exc_info.value.completed_result
    assert isinstance(completed, str)
    paths = _visual_paths(completed)
    assert len(paths) == 1
    assert _session_artifact_path(session_dir, paths[0]).exists()
    assert FilesystemTools(runtime, session_dir=session_dir).view_image(paths[0])[0].media_type == "image/png"


async def test_failed_conversion_cleans_new_image_orphans(tmp_path: Path) -> None:
    image = _png_bytes()

    class _FailingParser(_FakeParser):
        def parse(self, path: str, *, image_sink: DocumentImageSink | None = None) -> ParsedDocument:
            assert image_sink is not None
            assert image_sink.try_reserve_occurrence()
            assert image_sink.save_image(image, location="Page 1", ordinal=1, source_name="pixel.png") is not None
            raise RuntimeError("parse failed after image write")

    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    session_dir = tmp_path / "session"
    tools = DocConverterTools(_make_runtime(tmp_path), session_dir=session_dir)
    tools.set_image_extraction_enabled(True)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=_FailingParser()):
        result = await tools.convert_document(str(doc))

    assert result.startswith("Error:")
    assert list((session_dir / "doc_converter").glob("*.png")) == []


async def test_successful_conversion_cleans_images_missing_from_returned_occurrences(tmp_path: Path) -> None:
    image = _png_bytes()

    class _DroppingParser(_FakeParser):
        def parse(self, path: str, *, image_sink: DocumentImageSink | None = None) -> ParsedDocument:
            assert image_sink is not None
            assert image_sink.try_reserve_occurrence()
            assert image_sink.save_image(image, location="Page 1", ordinal=1, source_name="pixel.png") is not None
            return ParsedDocument(markdown="# Report")

    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")
    session_dir = tmp_path / "session"
    tools = DocConverterTools(_make_runtime(tmp_path), session_dir=session_dir)
    tools.set_image_extraction_enabled(True)

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=_DroppingParser()):
        result = await tools.convert_document(str(doc))

    assert "view_image" not in result
    assert list((session_dir / "doc_converter").glob("*.png")) == []


# ---------------------------------------------------------------------------
# Tool registration
# ---------------------------------------------------------------------------


def test_tools_returns_convert_document(tmp_path: Path) -> None:
    runtime = _make_runtime(tmp_path)
    dc = DocConverterTools(runtime, session_id="s1")
    tools = dc.tools()
    assert len(tools) == 1
    assert tools[0].name == "convert_document"
    assert get_tool_kind(tools[0]) == "doc_converter"
    assert tools[0].kind is None


# ---------------------------------------------------------------------------
# File collision handling
# ---------------------------------------------------------------------------


async def test_save_handles_filename_collision(tmp_path: Path, monkeypatch) -> None:
    """Second conversion of same-named file creates a _1 suffixed file."""
    monkeypatch.delenv(SESSION_ROOT_DIR_ENV_VAR, raising=False)
    doc = tmp_path / "report.pdf"
    doc.write_bytes(b"%PDF-fake")

    md_text = "# Report\n" + "content " * 100
    fake_parser = _FakeParser(md_text)

    runtime = _make_runtime(tmp_path)
    session_id = "collision_test"
    tools = DocConverterTools(runtime, session_id=session_id)

    # Pre-create the first output file to simulate collision
    out_dir = tmp_path / "sessions" / session_short_id(session_id) / "doc_converter"
    out_dir.mkdir(parents=True)
    (out_dir / "report.md").write_text("old content")

    with patch(f"{_PATCH_REGISTRY}.get_parser", return_value=fake_parser), patch(_PATCH_TOK) as mock_tok:
        mock_tok.count_tokens.return_value = _TOKEN_THRESHOLD + 100
        result = await tools.convert_document(str(doc))

    assert "report_1.md" in result
    assert (out_dir / "report_1.md").exists()


# ---------------------------------------------------------------------------
# Parser registry
# ---------------------------------------------------------------------------


def test_parser_registry_all_extensions() -> None:
    """All expected extensions are registered."""
    from chrys.service.tools.builtins.doc_converter.registry import supported_extensions

    exts = supported_extensions()
    assert ".pdf" in exts
    assert ".docx" in exts
    assert ".pptx" in exts
    assert ".xlsx" in exts
    assert ".xls" in exts
    assert ".epub" not in exts


def test_parser_registry_get_parser() -> None:
    """get_parser returns correct parser type for each extension."""
    from chrys.service.tools.builtins.doc_converter.parsers.docx import DocxParser
    from chrys.service.tools.builtins.doc_converter.parsers.pdf import PdfParser
    from chrys.service.tools.builtins.doc_converter.parsers.pptx import PptxParser
    from chrys.service.tools.builtins.doc_converter.parsers.xls import XlsParser
    from chrys.service.tools.builtins.doc_converter.parsers.xlsx import XlsxParser
    from chrys.service.tools.builtins.doc_converter.registry import get_parser

    assert isinstance(get_parser(".pdf"), PdfParser)
    assert isinstance(get_parser(".docx"), DocxParser)
    assert isinstance(get_parser(".pptx"), PptxParser)
    assert isinstance(get_parser(".xlsx"), XlsxParser)
    assert isinstance(get_parser(".xls"), XlsParser)
    assert get_parser(".epub") is None
    assert get_parser(".txt") is None


def test_parser_registry_case_insensitive() -> None:
    from chrys.service.tools.builtins.doc_converter.registry import get_parser

    assert get_parser(".PDF") is not None
    assert get_parser(".Docx") is not None


# ---------------------------------------------------------------------------
# Individual parser unit tests (mock the library imports)
# ---------------------------------------------------------------------------


def test_pdf_parser_output() -> None:
    """PdfParser produces page-based Markdown headings."""
    from chrys.service.tools.builtins.doc_converter.parsers.pdf import PdfParser

    mock_page1 = MagicMock()
    mock_page1.extract_text.return_value = "Hello world"
    mock_page2 = MagicMock()
    mock_page2.extract_text.return_value = "Second page content"

    mock_reader = MagicMock()
    mock_reader.pages = [mock_page1, mock_page2]

    mock_pypdf = MagicMock()
    mock_pypdf.PdfReader.return_value = mock_reader

    with patch.dict("sys.modules", {"pypdf": mock_pypdf}):
        parser = PdfParser()
        result = parser.parse("/fake.pdf")

    assert "# Page 1" in result.markdown
    assert "Hello world" in result.markdown
    assert "# Page 2" in result.markdown
    assert "Second page content" in result.markdown


def test_pdf_parser_extracts_image_only_page_and_keeps_no_text_placeholder(tmp_path: Path) -> None:
    from PIL import Image

    from chrys.service.tools.builtins.doc_converter.parsers.pdf import PdfParser

    pdf = tmp_path / "image.pdf"
    Image.new("RGB", (8, 8), (255, 0, 0)).save(pdf, format="PDF")
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="image")

    result = PdfParser().parse(str(pdf), image_sink=sink)

    assert "# Page 1\n\n(no text content)" in result.markdown
    assert len(result.visuals) == 1
    assert _session_artifact_path(tmp_path / "session", result.visuals[0].reference).exists()


def test_pdf_parser_deduplicates_repeated_indirect_reference(tmp_path: Path) -> None:
    from chrys.service.tools.builtins.doc_converter.parsers.pdf import PdfParser

    image_data = _png_bytes()
    indirect_reference = object()

    class _Images:
        def keys(self) -> list[str]:
            return ["/Im0"]

        def __getitem__(self, key: str):
            return MagicMock(
                data=image_data,
                name="logo.png",
                is_displayed=True,
                indirect_reference=indirect_reference,
            )

    pages = []
    for _ in range(2):
        page = MagicMock()
        page.extract_text.return_value = "text"
        page.images = _Images()
        pages.append(page)
    mock_pypdf = MagicMock()
    mock_pypdf.PdfReader.return_value.pages = pages
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="report")

    with (
        patch.dict("sys.modules", {"pypdf": mock_pypdf}),
        patch.object(sink, "save_image", wraps=sink.save_image) as save_image,
    ):
        result = PdfParser().parse("/fake.pdf", image_sink=sink)

    assert len(result.visuals) == 2
    assert result.visuals[0].reference == result.visuals[1].reference
    assert save_image.call_count == 1
    assert len(list((tmp_path / "session" / "doc_converter").iterdir())) == 1


def test_pdf_parser_skips_undisplayed_resource_without_writing(tmp_path: Path) -> None:
    from chrys.service.tools.builtins.doc_converter.parsers.pdf import PdfParser

    class _Images:
        def keys(self) -> list[str]:
            return ["/Unused"]

        def __getitem__(self, key: str):
            return MagicMock(
                data=_png_bytes(),
                name="unused.png",
                is_displayed=False,
                indirect_reference=object(),
            )

    page = MagicMock()
    page.extract_text.return_value = "text"
    page.images = _Images()
    mock_pypdf = MagicMock()
    mock_pypdf.PdfReader.return_value.pages = [page]
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="report")

    with patch.dict("sys.modules", {"pypdf": mock_pypdf}):
        result = PdfParser().parse("/fake.pdf", image_sink=sink)

    assert result.visuals == ()
    assert not (tmp_path / "session" / "doc_converter").exists()


def test_pdf_parser_corrupt_image_preserves_text_and_other_images(tmp_path: Path) -> None:
    from chrys.service.tools.builtins.doc_converter.parsers.pdf import PdfParser

    class _Images:
        def keys(self) -> list[str]:
            return ["/Corrupt", "/Valid"]

        def __getitem__(self, key: str):
            data = b"not-an-image" if key == "/Corrupt" else _png_bytes()
            return MagicMock(
                data=data,
                name=f"{key}.png",
                is_displayed=True,
                indirect_reference=None,
            )

    page = MagicMock()
    page.extract_text.return_value = "Preserved page text"
    page.images = _Images()
    mock_pypdf = MagicMock()
    mock_pypdf.PdfReader.return_value.pages = [page]
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="report")

    with patch.dict("sys.modules", {"pypdf": mock_pypdf}):
        result = PdfParser().parse("/fake.pdf", image_sink=sink)

    assert "Preserved page text" in result.markdown
    assert len(result.visuals) == 1
    assert len(result.warnings) == 1
    assert "could not be decoded or normalized" in result.warnings[0]


def test_pdf_parser_caps_explicit_image_retrievals(tmp_path: Path) -> None:
    from chrys.service.tools.builtins.doc_converter.parsers.pdf import PdfParser

    class _Images:
        def __init__(self) -> None:
            self.retrieval_count = 0

        def keys(self) -> list[str]:
            return [f"/Im{index}" for index in range(MAX_IMAGE_OCCURRENCES + 2)]

        def __getitem__(self, key: str):
            self.retrieval_count += 1
            return MagicMock(
                data=_png_bytes(),
                name="logo.png",
                is_displayed=True,
                indirect_reference=None,
            )

    images = _Images()
    page = MagicMock()
    page.extract_text.return_value = "text"
    page.images = images
    mock_pypdf = MagicMock()
    mock_pypdf.PdfReader.return_value.pages = [page]
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="report")

    with patch.dict("sys.modules", {"pypdf": mock_pypdf}):
        result = PdfParser().parse("/fake.pdf", image_sink=sink)

    assert images.retrieval_count == MAX_IMAGE_OCCURRENCES
    assert len(result.visuals) == MAX_IMAGE_OCCURRENCES
    assert any("2 image candidate(s)" in warning for warning in result.warnings)


def test_pdf_parser_documents_inline_scan_before_retrieval_budget(tmp_path: Path) -> None:
    from chrys.service.tools.builtins.doc_converter.parsers.pdf import PdfParser

    class _InlineImages:
        def __init__(self) -> None:
            self.inline_materialized = 0
            self.retrieval_count = 0

        def keys(self) -> list[str]:
            candidate_count = MAX_IMAGE_OCCURRENCES + 2
            self.inline_materialized += candidate_count
            return [f"~{index}~" for index in range(candidate_count)]

        def __getitem__(self, key: str):
            self.retrieval_count += 1
            return MagicMock(
                data=_png_bytes(),
                name="inline.png",
                is_displayed=True,
                indirect_reference=None,
            )

    images = _InlineImages()
    page = MagicMock()
    page.extract_text.return_value = "text"
    page.images = images
    mock_pypdf = MagicMock()
    mock_pypdf.PdfReader.return_value.pages = [page]
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="inline")

    with patch.dict("sys.modules", {"pypdf": mock_pypdf}):
        PdfParser().parse("/fake.pdf", image_sink=sink)

    assert images.inline_materialized == MAX_IMAGE_OCCURRENCES + 2
    assert images.retrieval_count == MAX_IMAGE_OCCURRENCES
    assert len(sink.warnings) == 1
    assert "2 image candidate(s)" in sink.warnings[0]


def test_docx_parser_headings_and_tables() -> None:
    """DocxParser maps heading styles to Markdown headings and extracts tables in order."""
    from chrys.service.tools.builtins.doc_converter.parsers.docx import DocxParser

    # Real classes so isinstance() checks work with the mocked docx.table.Table
    class _FakeDocxTable:
        def __init__(self, row_data: list[list[str]]):
            self.rows = [MagicMock(cells=[MagicMock(text=c) for c in r]) for r in row_data]

    mock_para1 = MagicMock()
    mock_para1.text = "My Title"
    mock_para1.style.name = "Heading 1"

    mock_para2 = MagicMock()
    mock_para2.text = "Some body text"
    mock_para2.style.name = "Normal"

    mock_table = _FakeDocxTable([["Name", "Age"], ["Alice", "30"]])

    mock_para3 = MagicMock()
    mock_para3.text = "Subsection"
    mock_para3.style.name = "Heading 2"

    mock_doc = MagicMock()
    mock_doc.iter_inner_content.return_value = iter([mock_para1, mock_para2, mock_table, mock_para3])

    mock_docx = MagicMock()
    mock_docx.Document.return_value = mock_doc
    mock_table_mod = MagicMock()
    mock_table_mod.Table = _FakeDocxTable

    with patch.dict("sys.modules", {"docx": mock_docx, "docx.table": mock_table_mod}):
        parser = DocxParser()
        result = parser.parse("/fake.docx")

    assert "# My Title" in result.markdown
    assert "Some body text" in result.markdown
    assert "| Name | Age |" in result.markdown
    assert "| Alice | 30 |" in result.markdown
    assert "## Subsection" in result.markdown
    # Verify ordering: title before table before subsection
    assert result.markdown.index("My Title") < result.markdown.index("Name | Age") < result.markdown.index("Subsection")


def test_docx_parser_extracts_package_wide_header_image(tmp_path: Path) -> None:
    from docx import Document

    from chrys.service.tools.builtins.doc_converter.parsers.docx import DocxParser

    document = Document()
    document.add_paragraph("Body text")
    document.sections[0].header.paragraphs[0].add_run().add_picture(BytesIO(_png_bytes()))
    path = tmp_path / "header.docx"
    document.save(path)
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="header")

    result = DocxParser().parse(str(path), image_sink=sink)

    assert result.markdown == "Body text"
    assert len(result.visuals) == 1
    assert result.visuals[0].location == "Document"
    assert _session_artifact_path(tmp_path / "session", result.visuals[0].reference).exists()


def test_docx_parser_package_image_parts_deduplicate_reused_blob(tmp_path: Path) -> None:
    from docx import Document

    from chrys.service.tools.builtins.doc_converter.parsers.docx import DocxParser

    image = BytesIO(_png_bytes())
    document = Document()
    document.add_picture(image)
    image.seek(0)
    document.add_picture(image)
    path = tmp_path / "reused.docx"
    document.save(path)
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="reused")

    result = DocxParser().parse(str(path), image_sink=sink)

    assert result.markdown == ""
    assert len(result.visuals) == 1
    assert len(list((tmp_path / "session" / "doc_converter").iterdir())) == 1


def test_pptx_parser_slides() -> None:
    """PptxParser produces slide-based Markdown headings."""
    from chrys.service.tools.builtins.doc_converter.parsers.pptx import PptxParser

    mock_title = MagicMock()
    mock_title.text = "Intro Slide"

    mock_shape = MagicMock()
    mock_shape.has_text_frame = True
    mock_shape.has_table = False
    mock_para = MagicMock()
    mock_para.text = "Bullet point"
    mock_shape.text_frame.paragraphs = [mock_para]

    mock_slide = MagicMock()
    mock_slide.shapes.title = mock_title
    mock_slide.shapes.__iter__ = MagicMock(return_value=iter([mock_shape]))

    mock_prs = MagicMock()
    mock_prs.slides = [mock_slide]

    mock_pptx = MagicMock()
    mock_pptx.Presentation.return_value = mock_prs

    with patch.dict("sys.modules", {"pptx": mock_pptx}):
        parser = PptxParser()
        result = parser.parse("/fake.pptx")

    assert "# Slide 1: Intro Slide" in result.markdown
    assert "Bullet point" in result.markdown


def test_pptx_parser_table_shape() -> None:
    """PptxParser extracts tables from table shapes."""
    from chrys.service.tools.builtins.doc_converter.parsers.pptx import PptxParser

    # Text shape
    text_shape = MagicMock()
    text_shape.has_table = False
    text_shape.has_text_frame = True
    text_shape.text_frame.paragraphs = [MagicMock(text="Intro text")]

    # Table shape
    table_shape = MagicMock()
    table_shape.has_table = True
    table_shape.has_text_frame = False
    table_shape.table.rows = [
        MagicMock(cells=[MagicMock(text="Col A"), MagicMock(text="Col B")]),
        MagicMock(cells=[MagicMock(text="val1"), MagicMock(text="val2")]),
    ]

    mock_slide = MagicMock()
    mock_slide.shapes.title = None
    mock_slide.shapes.__iter__ = MagicMock(return_value=iter([text_shape, table_shape]))

    mock_prs = MagicMock()
    mock_prs.slides = [mock_slide]

    mock_pptx = MagicMock()
    mock_pptx.Presentation.return_value = mock_prs

    with patch.dict("sys.modules", {"pptx": mock_pptx}):
        parser = PptxParser()
        result = parser.parse("/fake.pptx")

    assert "# Slide 1" in result.markdown
    assert "Intro text" in result.markdown
    assert "| Col A | Col B |" in result.markdown
    assert "| val1 | val2 |" in result.markdown


def test_pptx_parser_extracts_top_level_and_grouped_pictures(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    from chrys.service.tools.builtins.doc_converter.parsers.pptx import PptxParser

    image = tmp_path / "pixel.png"
    image.write_bytes(_png_bytes())
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    top_level = slide.shapes.add_picture(str(image), Inches(1), Inches(1))
    slide.shapes.add_group_shape([top_level])
    deck = tmp_path / "deck.pptx"
    presentation.save(deck)
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="deck")

    result = PptxParser().parse(str(deck), image_sink=sink)

    assert len(result.visuals) == 1
    assert result.visuals[0].location == "Slide 1"
    assert _session_artifact_path(tmp_path / "session", result.visuals[0].reference).exists()


def test_pptx_parser_extracts_populated_picture_placeholder(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.shapes.picture import Picture

    from chrys.service.tools.builtins.doc_converter.parsers.pptx import PptxParser

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    placeholder = next(
        candidate for candidate in slide.placeholders if candidate.placeholder_format.type == PP_PLACEHOLDER.PICTURE
    )
    populated = placeholder.insert_picture(BytesIO(_png_bytes()))
    assert isinstance(populated, Picture)
    assert populated.shape_type.name == "PLACEHOLDER"
    deck = tmp_path / "placeholder.pptx"
    presentation.save(deck)
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="placeholder")

    result = PptxParser().parse(str(deck), image_sink=sink)

    assert len(result.visuals) == 1
    assert _session_artifact_path(tmp_path / "session", result.visuals[0].reference).exists()


def test_pptx_parser_repeated_logo_across_slides_stores_once(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    from chrys.service.tools.builtins.doc_converter.parsers.pptx import PptxParser

    image = tmp_path / "logo.png"
    image.write_bytes(_png_bytes())
    presentation = Presentation()
    for _ in range(2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(image), Inches(1), Inches(1))
    deck = tmp_path / "repeated-logo.pptx"
    presentation.save(deck)
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="repeated-logo")

    result = PptxParser().parse(str(deck), image_sink=sink)

    assert len(result.visuals) == 2
    assert result.visuals[0].reference == result.visuals[1].reference
    assert len(list((tmp_path / "session" / "doc_converter").iterdir())) == 1


def test_pptx_parser_skips_linked_only_picture_with_bounded_warning(tmp_path: Path) -> None:
    from pptx.shapes.picture import Picture

    from chrys.service.tools.builtins.doc_converter.parsers.pptx import PptxParser

    class _LinkedPicture(Picture):
        @property
        def has_table(self) -> bool:
            return False

        @property
        def has_text_frame(self) -> bool:
            return False

        @property
        def image(self):
            raise ValueError("no embedded image")

    linked = object.__new__(_LinkedPicture)
    slide = MagicMock()
    slide.shapes.title = None
    slide.shapes.__iter__ = MagicMock(side_effect=lambda: iter([linked]))
    presentation = MagicMock()
    presentation.slides = [slide]
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="linked")

    with patch("pptx.Presentation", return_value=presentation):
        result = PptxParser().parse("/fake.pptx", image_sink=sink)

    assert result.visuals == ()
    assert len(result.warnings) == 1
    assert "linked image(s)" in result.warnings[0]


def test_xlsx_parser_tables() -> None:
    """XlsxParser produces sheet-based Markdown tables."""
    from chrys.service.tools.builtins.doc_converter.parsers.xlsx import XlsxParser

    mock_ws = MagicMock()
    mock_ws.iter_rows.return_value = [("Name", "Age"), ("Alice", 30), ("Bob", 25)]

    mock_wb = MagicMock()
    mock_wb.sheetnames = ["Sheet1"]
    mock_wb.__getitem__ = lambda self, key: mock_ws

    mock_openpyxl = MagicMock()
    mock_openpyxl.load_workbook.return_value = mock_wb

    with patch.dict("sys.modules", {"openpyxl": mock_openpyxl}):
        parser = XlsxParser()
        result = parser.parse("/fake.xlsx", image_sink=MagicMock(spec=DocumentImageSink))

    assert "# Sheet: Sheet1" in result.markdown
    assert "| Name | Age |" in result.markdown
    assert "| Alice | 30 |" in result.markdown
    assert result.warnings == ()


def test_xlsx_parser_with_embedded_image_stays_text_only_without_result_warning(tmp_path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as SpreadsheetImage

    from chrys.service.tools.builtins.doc_converter.parsers.xlsx import XlsxParser

    image_path = tmp_path / "chart.png"
    image_path.write_bytes(_png_bytes())
    workbook = Workbook()
    worksheet = workbook.active
    worksheet["A1"] = "Report"
    worksheet.add_image(SpreadsheetImage(image_path), "B2")
    spreadsheet = tmp_path / "report.xlsx"
    workbook.save(spreadsheet)
    workbook.close()
    sink = DocumentImageSink(tmp_path / "session" / "doc_converter", source_stem="report")

    result = XlsxParser().parse(str(spreadsheet), image_sink=sink)

    assert result.visuals == ()
    assert result.warnings == ()
    assert not (tmp_path / "session" / "doc_converter").exists()


def test_xls_parser_tables() -> None:
    """XlsParser produces sheet-based Markdown tables."""
    from chrys.service.tools.builtins.doc_converter.parsers.xls import XlsParser

    mock_sheet = MagicMock()
    mock_sheet.name = "Data"
    mock_sheet.nrows = 2
    mock_sheet.ncols = 2
    mock_sheet.cell_value = lambda r, c: [["Header1", "Header2"], ["val1", "val2"]][r][c]

    mock_wb = MagicMock()
    mock_wb.sheets.return_value = [mock_sheet]

    mock_xlrd = MagicMock()
    mock_xlrd.open_workbook.return_value = mock_wb

    with patch.dict("sys.modules", {"xlrd": mock_xlrd}):
        parser = XlsParser()
        result = parser.parse("/fake.xls", image_sink=MagicMock(spec=DocumentImageSink))

    assert "# Sheet: Data" in result.markdown
    assert "| Header1 | Header2 |" in result.markdown
    assert "| val1 | val2 |" in result.markdown
    assert result.warnings == ()


def test_protocol_compliance() -> None:
    """All parsers satisfy the DocParser protocol."""
    from chrys.service.tools.builtins.doc_converter.parsers.base import DocParser
    from chrys.service.tools.builtins.doc_converter.parsers.docx import DocxParser
    from chrys.service.tools.builtins.doc_converter.parsers.pdf import PdfParser
    from chrys.service.tools.builtins.doc_converter.parsers.pptx import PptxParser
    from chrys.service.tools.builtins.doc_converter.parsers.xls import XlsParser
    from chrys.service.tools.builtins.doc_converter.parsers.xlsx import XlsxParser

    for cls in [PdfParser, DocxParser, PptxParser, XlsxParser, XlsParser]:
        assert isinstance(cls(), DocParser), f"{cls.__name__} does not satisfy DocParser protocol"


# ---------------------------------------------------------------------------
# Markdown table pipe escaping
# ---------------------------------------------------------------------------


def test_table_pipe_escaping() -> None:
    """Pipe characters in cell values are escaped so they don't break the table."""
    from chrys.service.tools.builtins.doc_converter.parsers._table import rows_to_markdown_table

    rows = [("Header", "Formula"), ("A|B", "x | y")]
    result = rows_to_markdown_table(rows)
    assert "| Header | Formula |" in result
    assert r"| A\|B | x \| y |" in result
